"""Office 文档解析器

支持两种后端：
1. Unstructured — 高质量结构化解析（保留标题层级、表格结构）
2. 原生解析器 — 轻量级回退方案（python-docx / openpyxl / python-pptx）
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List

from doc_analyzer.contract import ParseResult, error_result
from doc_analyzer.model_config import get_model_config
from doc_analyzer.parsers.base import BaseParser

logger = logging.getLogger(__name__)


class OfficeParser(BaseParser):
    """Office 文档解析器（docx / xlsx / pptx）"""

    parser_name = "office"
    supported_extensions = [".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"]

    def __init__(self):
        self.unstructured_available = self._check_unstructured()

    def _check_unstructured(self) -> bool:
        """检查 Unstructured 是否可用"""
        try:
            from unstructured.partition.docx import partition_docx
            return True
        except ImportError:
            return False

    def parse(self, file_path: Path) -> ParseResult:
        file_path = Path(file_path).resolve()
        ext = file_path.suffix.lower()

        # .xls 无 Unstructured 支持，直接走原生（xlrd）
        if ext == ".xls":
            return self._parse_with_native(file_path, ext)
        # .doc / .ppt 优先走 Unstructured，失败再 fallback
        if ext in (".doc", ".ppt"):
            if self.unstructured_available:
                try:
                    return self._parse_with_unstructured(file_path, ext)
                except Exception as e:
                    logger.warning(f"Unstructured 解析 {ext} 失败，回退到原生解析器: {e}")
            return self._parse_with_native(file_path, ext)

        config = get_model_config("document", "office")
        provider = config.get("provider", "auto")

        if provider == "unstructured":
            if self.unstructured_available:
                return self._parse_with_unstructured(file_path, ext)
            else:
                logger.warning("Unstructured 不可用，回退到原生解析器")
                return self._parse_with_native(file_path, ext)
        elif provider == "native":
            return self._parse_with_native(file_path, ext)
        else:  # auto
            if self.unstructured_available:
                try:
                    return self._parse_with_unstructured(file_path, ext)
                except Exception as e:
                    logger.warning(f"Unstructured 解析失败，回退到原生解析器: {e}")
            return self._parse_with_native(file_path, ext)

        return self._parse_with_native(file_path, ext)

    # ========================================================================
    # Unstructured 后端
    # ========================================================================

    def _parse_with_unstructured(self, file_path: Path, ext: str) -> ParseResult:
        """使用 Unstructured 解析"""
        from doc_analyzer.parsers.unstructured_parser import UnstructuredOfficeParser

        parser = UnstructuredOfficeParser()

        if ext in (".docx", ".doc"):
            result = parser.parse_docx(file_path) if ext == ".docx" else parser.parse_doc(file_path)
        elif ext in (".xlsx", ".xls"):
            result = parser.parse_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            result = parser.parse_pptx(file_path) if ext == ".pptx" else parser.parse_ppt(file_path)
        else:
            raise ValueError(f"Unstructured 不支持: {ext}")

        metadata = {
            "filename": file_path.name,
            "file_size": file_path.stat().st_size,
            "parser_backend": "unstructured",
            "tables": len(result.tables),
            "sections": len(result.sections),
            "images": len(result.images),
            "page_count": result.page_count,
        }
        metadata.update(result.metadata)

        return ParseResult(
            source=str(file_path),
            type="document",
            format=ext.lstrip("."),
            content=result.markdown,
            metadata=metadata,
            parser=f"office(unstructured)",
        )

    # ========================================================================
    # 原生解析器回退
    # ========================================================================

    def _parse_with_native(self, file_path: Path, ext: str) -> ParseResult:
        """使用原生解析器"""
        if ext in (".docx", ".doc"):
            return self._parse_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return self._parse_xlsx(file_path)
        elif ext in (".pptx", ".ppt"):
            return self._parse_pptx(file_path)
        else:
            return error_result(str(file_path), f"不支持的格式: {ext}", ext.lstrip("."))

    def _parse_docx(self, file_path: Path) -> ParseResult:
        """解析 Word 文档（原生）"""
        try:
            from docx import Document

            doc = Document(str(file_path))
            text_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    # 根据样式判断标题级别
                    style_name = para.style.name if para.style else ""
                    if style_name.startswith("Heading"):
                        try:
                            level = int(style_name.replace("Heading ", "").strip())
                            prefix = "#" * min(level, 6)
                            text_parts.append(f"\n{prefix} {para.text.strip()}\n")
                        except ValueError:
                            text_parts.append(para.text.strip())
                    else:
                        text_parts.append(para.text.strip())

            # 提取表格
            for i, table in enumerate(doc.tables, 1):
                text_parts.append(f"\n[表格 {i}]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_parts.append(row_text)

            content = "\n\n".join(text_parts)

            if not content.strip():
                return error_result(str(file_path), "Word 文档无文本内容", "docx")

            metadata = {
                "filename": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser_backend": "python-docx",
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            }

            return ParseResult(
                source=str(file_path),
                type="document",
                format="docx",
                content=content,
                metadata=metadata,
                parser="office(native)",
            )
        except Exception as e:
            return error_result(str(file_path), f"Word 解析失败: {e}", "docx")

    def _parse_xlsx(self, file_path: Path) -> ParseResult:
        """解析 Excel 文档（原生）"""
        try:
            ext = file_path.suffix.lower()

            if ext == ".xls":
                # 老格式用 xlrd
                import xlrd
                wb = xlrd.open_workbook(str(file_path))
                sheets = wb.sheets()
                sheet_names = [s.name for s in sheets]
                text_parts = []
                for sheet in sheets:
                    text_parts.append(f"# 工作表: {sheet.name}\n")
                    for row_idx in range(sheet.nrows):
                        row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols) if sheet.cell_value(row_idx, col) != ""]
                        row_text = " | ".join(row)
                        if row_text.strip():
                            text_parts.append(row_text)
                content_text = "\n".join(text_parts)
                if not content_text.strip():
                    return error_result(str(file_path), "Excel 无文本内容", "xls")
                metadata = {
                    "filename": file_path.name,
                    "file_size": file_path.stat().st_size,
                    "parser_backend": "xlrd",
                    "sheets": len(sheet_names),
                    "sheet_names": sheet_names,
                }
                return ParseResult(
                    source=str(file_path),
                    type="document",
                    format="xls",
                    content=content_text,
                    metadata=metadata,
                    parser="office(native/xlrd)",
                )

            # .xlsx 用 openpyxl
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"# 工作表: {sheet_name}\n")

                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(
                        str(cell) for cell in row if cell is not None
                    )
                    if row_text.strip():
                        text_parts.append(row_text)

            content = "\n".join(text_parts)

            if not content.strip():
                return error_result(str(file_path), "Excel 无文本内容", "xlsx")

            metadata = {
                "filename": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser_backend": "openpyxl",
                "sheets": len(wb.sheetnames),
                "sheet_names": wb.sheetnames,
            }

            return ParseResult(
                source=str(file_path),
                type="document",
                format="xlsx",
                content=content,
                metadata=metadata,
                parser="office(native)",
            )
        except Exception as e:
            return error_result(str(file_path), f"Excel 解析失败: {e}", "xlsx")

    def _parse_pptx(self, file_path: Path) -> ParseResult:
        """解析 PowerPoint 文档（原生）"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                text_parts.append(f"# 幻灯片 {slide_num}\n")

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text.strip())

            content = "\n\n".join(text_parts)

            if not content.strip():
                return error_result(str(file_path), "PowerPoint 无文本内容", "pptx")

            metadata = {
                "filename": file_path.name,
                "file_size": file_path.stat().st_size,
                "parser_backend": "python-pptx",
                "slides": len(prs.slides),
            }

            return ParseResult(
                source=str(file_path),
                type="document",
                format="pptx",
                content=content,
                metadata=metadata,
                parser="office(native)",
            )
        except Exception as e:
            return error_result(str(file_path), f"PowerPoint 解析失败: {e}", "pptx")
